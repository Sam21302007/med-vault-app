import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# File Paths
IMG_HERO = r"C:\Users\ASUS\.gemini\antigravity\brain\e40a6f96-5dd7-4e9b-9ed8-ee5a34b4c813\medvault_dashboard_hero_1786606988473.png"
IMG_ARCH = r"C:\Users\ASUS\.gemini\antigravity\brain\e40a6f96-5dd7-4e9b-9ed8-ee5a34b4c813\medvault_architecture_diagram_1786607004758.png"
IMG_PORTAL = r"C:\Users\ASUS\.gemini\antigravity\brain\e40a6f96-5dd7-4e9b-9ed8-ee5a34b4c813\medvault_patient_portal_1786607189308.png"

# Color Palettes
COLOR_DARK_BG = RGBColor(10, 15, 29)        # Deep Slate Navy
COLOR_DARK_CARD = RGBColor(30, 41, 59)      # Slate Card
COLOR_DARK_TEXT = RGBColor(241, 245, 249)    # Bright Off-White
COLOR_DARK_SUBTEXT = RGBColor(148, 163, 184) # Muted Silver

COLOR_LIGHT_BG = RGBColor(248, 250, 252)    # Soft Slate White
COLOR_LIGHT_CARD = RGBColor(255, 255, 255)  # Pure White Card
COLOR_LIGHT_TEXT = RGBColor(15, 23, 42)      # Deep Navy Text
COLOR_LIGHT_SUBTEXT = RGBColor(71, 85, 105)  # Slate Muted Text

COLOR_PRIMARY = RGBColor(6, 182, 212)       # Cyan Accent
COLOR_SECONDARY = RGBColor(59, 130, 246)    # Royal Blue Accent
COLOR_EMERALD = RGBColor(16, 185, 129)      # Emerald Green Accent

SLIDES_DATA = [
    {
        "title": "MedVault: Hospital Management System",
        "subtitle": "Next-Generation Digital Healthcare & Firestore EMR Platform",
        "bullets": [
            "Full-Stack Solution: Integrated React 18, Node.js Express, and Firebase Cloud",
            "Seamless Database Engine: Migrated from MongoDB to Firebase Firestore & Authentication",
            "Zero-Crash Reliability: Includes custom in-memory fallback for offline testing",
            "Role-Based Ecosystem: Personalized interfaces for Patients, Doctors, and Administrators"
        ],
        "image": IMG_HERO
    },
    {
        "title": "Executive Summary & Industry Solution",
        "subtitle": "Solving Data Fragmentation & Patient Care Delays",
        "bullets": [
            "Unified Medical Infrastructure: Centralizes EMR, Appointments, Pharmacy & Billing",
            "Real-Time Data Access: Instant sync between patient submissions and doctor terminals",
            "Paperless Operations: Digital prescriptions, electronic invoicing, and hospital audit trail",
            "Enterprise Security: Encrypted auth tokens, Firebase Security Rules, and RBAC control"
        ],
        "image": None
    },
    {
        "title": "System Architecture & Modern Tech Stack",
        "subtitle": "High-Performance Cloud Architecture Blueprint",
        "bullets": [
            "Frontend Layer: React 18, React Router 6, Tailwind CSS & Glassmorphism UI Design",
            "Backend Gateway: Node.js Express REST API server running port-resilient routing",
            "Database Engine: Firebase Firestore NoSQL document database with live collection sync",
            "Authentication: Firebase Authentication SDK & JSON Web Token (JWT) fallback"
        ],
        "image": IMG_ARCH
    },
    {
        "title": "Firebase Firestore Migration Story",
        "subtitle": "Transitioning from Legacy MongoDB to Firebase Cloud",
        "bullets": [
            "Database Modernization: Complete purge of Mongoose & MongoDB dependencies",
            "Firestore Collections: Structured models for users, appointments, beds, pharmacy, and billing",
            "Zero-Crash Fallback: Built-in InMemoryDb emulator for instant zero-config execution",
            "Live Cloud Sync: Seamless sync with Firebase Admin SDK key authentication"
        ],
        "image": None
    },
    {
        "title": "Multi-Role Ecosystem & User Governance",
        "subtitle": "Tailored Interfaces for All Stakeholders",
        "bullets": [
            "Patient Portal: Self-service booking, consultation history, and prescription downloads",
            "Doctor Terminal: Clinical note entry, appointment approval, and patient record review",
            "Admin Center: Hospital wide bed occupancy, inventory reordering, and audit logs",
            "Unified Access: Single login gateway with automated role routing"
        ],
        "image": IMG_PORTAL
    },
    {
        "title": "Patient Portal & Digital Experience",
        "subtitle": "Empowering Patients with Seamless Self-Service",
        "bullets": [
            "Instant Booking: Real-time appointment scheduling with priority tags (Urgent/Normal)",
            "Medical History: Comprehensive timeline of past diagnoses, prescriptions, and lab notes",
            "Digital Invoicing: Downloadable PDF billing receipts and online payment tracking",
            "Profile Management: Centralized personal health info, emergency contact, and vitals"
        ],
        "image": None
    },
    {
        "title": "Doctor Workstation & Clinical Workflows",
        "subtitle": "Optimized Tools for Healthcare Professionals",
        "bullets": [
            "Appointment Queue: Live dashboard displaying daily patient rosters and urgent cases",
            "EMR Authoring: Fast record creation for diagnosis, medication prescriptions, and follow-ups",
            "Availability Scheduler: Customizable weekly availability slots and consultation duration",
            "Patient Search: Rapid access to patient medical histories during consultations"
        ],
        "image": None
    },
    {
        "title": "Hospital Operations: Bed & Pharmacy Management",
        "subtitle": "Maximizing Resource Utilization & Supply Chains",
        "bullets": [
            "Live Bed Tracking: Real-time monitor for ICU, General Ward, and VIP Suite occupancy",
            "Admissions & Status: Dynamic allocation, status updates (Occupied/Cleaning/Available)",
            "Pharmacy Inventory: Automated stock tracking with instant low-stock alerts",
            "Batch Management: Expiry date monitoring and medicine restocking workflows"
        ],
        "image": None
    },
    {
        "title": "Security, Compliance & Immutable Auditing",
        "subtitle": "Protecting Patient Privacy & Regulatory Integrity",
        "bullets": [
            "Role-Based Access Control (RBAC): Strict endpoint protection enforcing Doctor/Patient scope",
            "Firebase Security: Enterprise-grade access control via Firebase Admin SDK",
            "System Audit Trail: Comprehensive log recording every bed assignment, invoice, and stock change",
            "Password Protection: Industry-standard bcrypt hashing for credential security"
        ],
        "image": None
    },
    {
        "title": "Deployment, Scalability & Strategic Roadmap",
        "subtitle": "Future-Proofing MedVault for Enterprise Scale",
        "bullets": [
            "Cloud Deployment: Production ready Vercel serverless hosting with global CDN",
            "CI/CD Integration: GitHub repository sync for automated zero-downtime builds",
            "AI Assistant Integration: Planned AI-assisted preliminary triage and diagnostics",
            "Telemedicine Expansion: Upcoming WebRTC video consultation & remote patient monitoring"
        ],
        "image": IMG_HERO
    }
]

def create_pptx(filename, is_dark=True):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg_color = COLOR_DARK_BG if is_dark else COLOR_LIGHT_BG
    card_color = COLOR_DARK_CARD if is_dark else COLOR_LIGHT_CARD
    text_color = COLOR_DARK_TEXT if is_dark else COLOR_LIGHT_TEXT
    subtext_color = COLOR_DARK_SUBTEXT if is_dark else COLOR_LIGHT_SUBTEXT

    for i, slide_data in enumerate(SLIDES_DATA):
        slide = prs.slides.add_slide(blank_layout)

        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

        # Top Accent Line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = COLOR_PRIMARY
        top_line.line.fill.background()

        # Header Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = text_color

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = slide_data["subtitle"]
        sp.font.size = Pt(16)
        sp.font.color.rgb = COLOR_PRIMARY

        # Determine Layout: With Image vs Bullets Only
        has_image = slide_data["image"] and os.path.exists(slide_data["image"])
        content_width = Inches(6.5) if has_image else Inches(11.7)

        # Bullet Cards Container / Box
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), content_width, Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.fill.background()

        bullet_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), content_width - Inches(0.4), Inches(4.6))
        btf = bullet_box.text_frame
        btf.word_wrap = True

        for b_idx, bullet_text in enumerate(slide_data["bullets"]):
            bp = btf.add_paragraph() if b_idx > 0 else btf.paragraphs[0]
            bp.text = f"•  {bullet_text}"
            bp.font.size = Pt(15)
            bp.font.color.rgb = subtext_color
            bp.space_after = Pt(14)

        # Insert Image if available
        if has_image:
            img_path = slide_data["image"]
            slide.shapes.add_picture(img_path, Inches(7.6), Inches(1.9), width=Inches(4.9))

        # Footer Slide Number
        footer_box = slide.shapes.add_textbox(Inches(11.0), Inches(7.0), Inches(1.5), Inches(0.4))
        fp = footer_box.text_frame.paragraphs[0]
        fp.text = f"Slide {i+1} of 10"
        fp.font.size = Pt(11)
        fp.font.color.rgb = subtext_color
        fp.alignment = PP_ALIGN.RIGHT

    prs.save(filename)
    print(f"[OK] Generated PowerPoint Presentation: {filename}")

def create_pdf(filename, is_dark=True):
    doc = SimpleDocTemplate(
        filename,
        pagesize=landscape(letter),
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    bg_color = colors.HexColor("#0A0F1D") if is_dark else colors.HexColor("#F8FAFC")
    card_color = colors.HexColor("#1E293B") if is_dark else colors.HexColor("#FFFFFF")
    text_color = colors.HexColor("#F1F5F9") if is_dark else colors.HexColor("#0F172A")
    subtext_color = colors.HexColor("#94A3B8") if is_dark else colors.HexColor("#475569")
    accent_color = colors.HexColor("#06B6D4")

    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=24,
        leading=28,
        textColor=text_color,
        spaceAfter=4
    )
    subtitle_style = ParagraphStyle(
        'SlideSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=accent_color,
        spaceAfter=15
    )
    bullet_style = ParagraphStyle(
        'SlideBullet',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=13,
        leading=18,
        textColor=subtext_color,
        spaceAfter=10
    )

    elements = []

    for i, slide_data in enumerate(SLIDES_DATA):
        elements.append(Paragraph(slide_data["title"], title_style))
        elements.append(Paragraph(slide_data["subtitle"], subtitle_style))

        # Bullet List Items
        bullet_pars = [Paragraph(f"• &nbsp; {b}", bullet_style) for b in slide_data["bullets"]]

        has_image = slide_data["image"] and os.path.exists(slide_data["image"])
        if has_image:
            img = RLImage(slide_data["image"], width=300, height=210)
            table_data = [[bullet_pars, img]]
            col_widths = [400, 310]
        else:
            table_data = [[bullet_pars]]
            col_widths = [710]

        t = Table(table_data, colWidths=col_widths)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), card_color),
            ('TOPPADDING', (0,0), (-1,-1), 16),
            ('BOTTOMPADDING', (0,0), (-1,-1), 16),
            ('LEFTPADDING', (0,0), (-1,-1), 16),
            ('RIGHTPADDING', (0,0), (-1,-1), 16),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#334155") if is_dark else colors.HexColor("#E2E8F0")),
        ]))

        elements.append(t)
        
        if i < len(SLIDES_DATA) - 1:
            elements.append(PageBreak())

    def draw_bg(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(bg_color)
        canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=True, stroke=False)
        canvas.setFillColor(accent_color)
        canvas.rect(0, doc.pagesize[1] - 8, doc.pagesize[0], 8, fill=True, stroke=False)
        canvas.restoreState()

    doc.build(elements, onFirstPage=draw_bg, onLaterPages=draw_bg)
    print(f"[OK] Generated PDF Presentation: {filename}")

if __name__ == "__main__":
    out_dir = r"d:\SVAC REACT"
    dark_pptx = os.path.join(out_dir, "MedVault_Presentation_Dark_Theme.pptx")
    light_pptx = os.path.join(out_dir, "MedVault_Presentation_Light_Theme.pptx")
    dark_pdf = os.path.join(out_dir, "MedVault_Presentation_Dark_Theme.pdf")
    light_pdf = os.path.join(out_dir, "MedVault_Presentation_Light_Theme.pdf")

    create_pptx(dark_pptx, is_dark=True)
    create_pptx(light_pptx, is_dark=False)
    create_pdf(dark_pdf, is_dark=True)
    create_pdf(light_pdf, is_dark=False)
